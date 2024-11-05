from flask import Flask, request, jsonify, send_file
import os
from werkzeug.utils import secure_filename
import pdfplumber
import groq
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Pt
import markdown2
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import io
import re

app = Flask(__name__)

app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024 # 50 MB max limit.

UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

CORS(app)

client = groq.Groq(api_key=os.environ["GROQ_API_KEY"])


@app.route('/')
def home():
    return "Welcome to the PDF Story App"


def extractText(file):
    text = ''
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text+= page.extract_text()
    return text

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)

    text = extractText(file_path)  # Assuming you have a function `extract_text` that extracts text from the PDF.
    story = generateStory(text)
    

    return jsonify({'message': 'File uploaded and story generated', 'story': story, 'file_path': file_path}), 200

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'pdf'}


def generateStory(text):
    messages = [
        {
            "role": "system",
            "content": "You are a helpful assistant, for a student, that generates a story, which is based upon the     content i provide. The PDF is read, and the text is given to you. You need to generate an entire story, which explains the given text content. This basically helps user to understand the content, like a quick summary. And most importantly, highlight the words to be remembered. Do not give questions. Only the Story. Make the story captivating, and also include emojis in the story. No need to mention them separately at the end."
        },
        {
            "role": "user",
            "content": text
        }
    ]
    
    # Call Groq's API to generate the chat completion
    chat_completion = client.chat.completions.create(
        messages=messages,
        model="llama-3.1-70b-versatile",  # You can select from various models like 'Gemma', 'Mixtral', etc.
        temperature=0.7,  # Adjust as needed for response randomness
        max_tokens=1024,  # Adjust token length for your needs
        top_p=1,  # Set nucleus sampling for diversity
        stream=False  # Set this to True for streaming responses
    )

    # Return the generated content
    return chat_completion.choices[0].message.content


def strip_html_tags(text):
    # Regular expression to remove HTML tags
    clean = re.compile('<.*?>')
    return re.sub(clean, '', text)

def convert_markdown_to_text_frame(text_frame, markdown_content):
    # Convert Markdown content to HTML, then remove tags and process for PowerPoint
    html_content = markdown2.markdown(markdown_content)
    
    # Split the HTML content and process each line
    for line in html_content.splitlines():
        plain_text = strip_html_tags(line)  # Remove any remaining HTML tags
        
        # Determine style based on HTML tags before stripping
        if line.startswith("<h1>"):
            p = text_frame.add_paragraph()
            p.text = plain_text
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        elif line.startswith("<h2>"):
            p = text_frame.add_paragraph()
            p.text = plain_text
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
        elif line.startswith("<ul>"):
            items = line.split("<li>")
            for item in items[1:]:
                item_text = strip_html_tags(item.split("</li>")[0])
                p = text_frame.add_paragraph()
                p.text = item_text
                p.level = 1  # Bullet point level
                p.font.size = Pt(20)
        elif "<strong>" in line:
            p = text_frame.add_paragraph()
            p.text = plain_text
            p.font.size = Pt(20)
            p.font.bold = True
        else:
            p = text_frame.add_paragraph()
            p.text = plain_text
            p.font.size = Pt(20)


@app.route('/download_ppt', methods=['POST'])
def download_ppt():
    paragraphs = request.json.get('paragraphs', [])
    ppt = Presentation()
    
    # Add a title slide
    title_slide_layout = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Flashcards Generated from PDF"
    pdf_name = request.json.get('pdf_name', 'Flashcards')
    
    # Add each paragraph as a slide with Markdown formatting
    for para in paragraphs:
        # Use a basic text slide layout
        slide_layout = ppt.slide_layouts[1]  # Title and Content layout
        slide = ppt.slides.add_slide(slide_layout)
        
        # Add title and content
        title = slide.shapes.title
        title.text = pdf_name
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()  # Clear any default content
        convert_markdown_to_text_frame(text_frame, para)  # Format and add the paragraph as Markdown

    # Save the presentation to an in-memory file for sending
    ppt_stream = io.BytesIO()
    ppt.save(ppt_stream)
    ppt_stream.seek(0)

    return send_file(ppt_stream, as_attachment=True, download_name=f"{pdf_name}.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

@app.route('/delete_file', methods=['POST'])
def delete_file():
    file_path = request.json.get('file_path')
    if file_path and os.path.exists(file_path):
        os.remove(file_path)
        return jsonify({"status": "File deleted successfully"})
    return jsonify({"status": "File not found"}), 404

if __name__ == "__main__":
    app.run(debug=True)


if __name__ == '__main__':
    app.run(debug=True)



